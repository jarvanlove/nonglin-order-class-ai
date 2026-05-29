"""
Export HTML slides to PPTX (one file per day).
Screenshots each HTML slide via Playwright, then assembles into separate .pptx files.
Structure: day01/index.html ~ day05/index.html, each with multiple slides.
Usage: python export_pptx.py
"""
import os
import sys
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

SLIDES_DIR = Path(__file__).parent.parent.resolve()
SCREENSHOT_DIR = Path(__file__).parent / "screenshots"
OUTPUT_DIR = SLIDES_DIR / ".." / "outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Clean old screenshots
SCREENSHOT_DIR.mkdir(exist_ok=True)
for old in SCREENSHOT_DIR.glob("*.png"):
    old.unlink()

DAYS = [f"day{i:02d}" for i in range(1, 6)]

VIEWPORT = {"width": 1280, "height": 720}
DEVICE_SCALE = 1


def screenshot_day(page, day_dir):
    """Screenshot all slides in a day's index.html. Returns list of screenshot paths."""
    src = SLIDES_DIR / day_dir / "index.html"
    if not src.exists():
        print(f"  [{day_dir}] File not found: {src}")
        return []

    page.goto(f"file:///{src.as_posix()}")
    page.wait_for_timeout(2000)

    # Ensure slides don't shrink in flex layout (fixes day02-day06 cover slide width)
    page.add_style_tag(content=".slide { flex-shrink: 0 !important; }")

    total = page.evaluate("""() => {
        const slides = document.querySelectorAll('.slide');
        return slides.length;
    }""")
    print(f"  [{day_dir}] Found {total} slides")

    results = []
    for i in range(total):
        page.evaluate(f"""(() => {{
            const slides = document.querySelectorAll('.slide');
            slides.forEach((s, idx) => s.classList.toggle('active', idx === {i}));
            // If active slide is outside deck-container, hide deck-container to prevent flex sharing
            const active = document.querySelector('.slide.active');
            const deck = document.querySelector('.deck-container');
            if (deck && active && !active.closest('.deck-container')) {{
                deck.style.display = 'none';
            }} else if (deck) {{
                deck.style.display = '';
            }}
        }})();""")
        page.wait_for_timeout(800)

        dst = SCREENSHOT_DIR / f"{day_dir}_slide{i+1:02d}.png"
        page.screenshot(path=str(dst))
        results.append(dst)

    return results


def build_pptx(shots, output_path):
    """Assemble screenshots into a PPTX file."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for shot in shots:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(shot),
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

    prs.save(str(output_path))
    print(f"  Saved: {output_path.name} ({len(shots)} slides)")


with sync_playwright() as p:
    browser = p.chromium.launch()
    context = browser.new_context(viewport=VIEWPORT, device_scale_factor=DEVICE_SCALE)
    page = context.new_page()

    for day in DAYS:
        day_path = SLIDES_DIR / day / "index.html"
        if not day_path.exists():
            print(f"[{day}] Skipping (no index.html)")
            continue
        print(f"[{day}] Screenshotting...")
        shots = screenshot_day(page, day)
        if not shots:
            continue

        # Build one PPTX per day
        output_path = OUTPUT_DIR / f"{day}_课件.pptx"
        build_pptx(shots, output_path)

    browser.close()

print(f"\nAll PPTX files saved to: {OUTPUT_DIR}")
